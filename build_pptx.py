"""
Build the Incluzza client pitch + roadmap PowerPoint deck.

Mirrors the website structure exactly:
  • Cover
  • Part 1: Client pitch (5 slides)
  • Part 2: Roadmap (5 slides)
  • Close

Design uses the same light & airy palette as the site:
  - surface: #FCFCFB
  - ink: #0B0B0B
  - accent: #2A78D6
  - categorical accents: #EB6834, #1BAF7A, #EDA100, #4A3AA7, #E87BA4, #008300, #E34948
  - serif: Calibri Light is used as a stand-in for Fraunces (PowerPoint has no native Fraunces)
  - sans:  Calibri
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Palette ────────────────────────────────────────────────────────
SURFACE       = RGBColor(0xFC, 0xFC, 0xFB)
SURFACE_2     = RGBColor(0xF5, 0xF3, 0xEC)
PAPER         = RGBColor(0xFF, 0xFF, 0xFF)
INK           = RGBColor(0x0B, 0x0B, 0x0B)
INK_2         = RGBColor(0x2A, 0x2A, 0x28)
INK_3         = RGBColor(0x52, 0x51, 0x4E)
MUTED         = RGBColor(0x89, 0x87, 0x81)
LINE          = RGBColor(0xE7, 0xE5, 0xDD)
ACCENT        = RGBColor(0x2A, 0x78, 0xD6)
ACCENT_SOFT   = RGBColor(0xD6, 0xE6, 0xFB)
ACCENT_WARM   = RGBColor(0xF4, 0xE3, 0xD6)
CAT = [
    RGBColor(0x2A, 0x78, 0xD6),
    RGBColor(0xEB, 0x68, 0x34),
    RGBColor(0x1B, 0xAF, 0x7A),
    RGBColor(0xED, 0xA1, 0x00),
    RGBColor(0x4A, 0x3A, 0xA7),
    RGBColor(0xE8, 0x7B, 0xA4),
    RGBColor(0x00, 0x83, 0x00),
    RGBColor(0xE3, 0x49, 0x48),
]

SERIF = "Calibri Light"
SANS  = "Calibri"

# ── Presentation setup (16:9) ─────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


# ── Helpers ────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w is not None:
            s.line.width = line_w
    s.text_frame.margin_left = s.text_frame.margin_right = 0
    s.text_frame.margin_top = s.text_frame.margin_bottom = 0
    return s


def add_text(slide, x, y, w, h, text, *,
             size=18, bold=False, italic=False, color=INK,
             font=SANS, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15, space_after=0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    # Support multi-line via newline split
    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_runs(slide, x, y, w, h, runs, *,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    """runs: list of (text, opts) where opts is a dict with size/bold/italic/color/font."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for text, opts in runs:
        r = p.add_run()
        r.text = text
        r.font.name = opts.get("font", SANS)
        r.font.size = Pt(opts.get("size", 16))
        r.font.bold = opts.get("bold", False)
        r.font.italic = opts.get("italic", False)
        r.font.color.rgb = opts.get("color", INK)
    return tb


def slide_background(slide, color=SURFACE):
    bg = add_rect(slide, 0, 0, SW, SH, fill=color, shape=MSO_SHAPE.RECTANGLE)
    # Send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def add_ornament_orbs(slide):
    """Soft pastel orbs (simulated with low-opacity rectangles — PowerPoint's
    blur is heavy; use 3 large soft shapes that read as ambient color)."""
    orbs = [
        (Inches(-1.0), Inches(-1.5), Inches(5.5), Inches(5.5), RGBColor(0xD6, 0xE6, 0xFB)),
        (Inches(10.0), Inches(0.5),  Inches(5.0), Inches(5.0), RGBColor(0xF4, 0xE3, 0xD6)),
        (Inches(2.0),  Inches(4.0),  Inches(6.0), Inches(6.0), RGBColor(0xD6, 0xF1, 0xE3)),
    ]
    for (x, y, w, h, color) in orbs:
        s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
        s.shadow.inherit = False
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.fill.transparency = 0.55
        s.line.fill.background()


def kicker_pill(slide, x, y, text, color=ACCENT):
    w = Inches(2.6)
    h = Inches(0.32)
    pill = add_rect(slide, x, y, w, h, fill=PAPER, line=LINE, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = pill.text_frame
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text.upper()
    r.font.name = SANS
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = color
    r.text._r.set("xml:space", "preserve")
    # set letter spacing via XML
    rPr = r._r.get_or_add_rPr()
    rPr.set("spc", "300")


# ── Slide builders ─────────────────────────────────────────────────
def slide_cover():
    s = prs.slides.add_slide(BLANK)
    slide_background(s)
    add_ornament_orbs(s)

    # Eyebrow
    add_text(s, Inches(0.8), Inches(1.4), Inches(8), Inches(0.4),
             "A partner for disability inclusion",
             size=14, color=INK_3, font=SANS)

    # Big headline
    add_runs(s, Inches(0.8), Inches(1.9), Inches(12), Inches(2.8),
             [("Disability inclusion,\n", {"size": 84, "font": SERIF, "color": INK}),
              ("made practical.", {"size": 84, "font": SERIF, "color": ACCENT, "italic": True})],
             line_spacing=1.0)

    # Lede
    add_text(s, Inches(0.8), Inches(4.7), Inches(10), Inches(1.2),
             "Incluzza helps your organization move from good intentions to a workplace and customer experience that genuinely works for people with disabilities.",
             size=18, color=INK_3, line_spacing=1.45)

    # Meta line
    add_runs(s, Inches(0.8), Inches(6.4), Inches(10), Inches(0.4),
             [("Prepared for your team   ·   ", {"size": 13, "color": MUTED}),
              ("20-minute introduction", {"size": 13, "color": MUTED})])

    # Brand mark in corner
    brand_x, brand_y = Inches(11.6), Inches(0.5)
    mark = s.shapes.add_shape(MSO_SHAPE.OVAL, brand_x, brand_y, Inches(0.5), Inches(0.5))
    mark.shadow.inherit = False
    mark.fill.solid()
    mark.fill.fore_color.rgb = ACCENT_SOFT
    mark.line.fill.background()
    add_text(s, brand_x, brand_y, Inches(0.5), Inches(0.5),
             "i", size=22, italic=True, font=SERIF, color=INK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, brand_x + Inches(0.6), brand_y, Inches(1.5), Inches(0.5),
             "Incluzza", size=15, bold=True, color=INK,
             anchor=MSO_ANCHOR.MIDDLE)


def slide_part_marker(part_num, title, sub):
    s = prs.slides.add_slide(BLANK)
    slide_background(s, SURFACE_2)
    add_ornament_orbs(s)

    # Top + bottom hairlines
    add_rect(s, Inches(2), Inches(2.6), Inches(9.33), Emu(9525), fill=LINE)
    add_rect(s, Inches(2), Inches(4.9), Inches(9.33), Emu(9525), fill=LINE)

    # Centered pill
    pill = add_rect(s, Inches(5.7), Inches(2.85), Inches(1.93), Inches(0.34),
                    fill=PAPER, line=ACCENT_SOFT, line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = pill.text_frame
    tf.margin_left = tf.margin_right = Inches(0.15)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = part_num.upper()
    r.font.name = SANS
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = ACCENT
    rPr = r._r.get_or_add_rPr()
    rPr.set("spc", "300")

    # Title
    add_text(s, Inches(1), Inches(3.5), Inches(11.33), Inches(1.0),
             title, size=56, font=SERIF, color=INK, align=PP_ALIGN.CENTER,
             line_spacing=1.05)

    # Subtitle
    add_text(s, Inches(1), Inches(4.4), Inches(11.33), Inches(0.5),
             sub, size=18, color=INK_3, align=PP_ALIGN.CENTER)


def slide_kicker_section(kicker, title_html_runs, sub=None):
    """A 'kicker + big title + sub' page header pattern. title_html_runs is a list of
    (text, opts) — exactly like add_runs — for the H2 line."""
    s = prs.slides.add_slide(BLANK)
    slide_background(s)
    add_ornament_orbs(s)
    kicker_pill(s, Inches(5.37), Inches(1.4), kicker)
    add_runs(s, Inches(1), Inches(2.0), Inches(11.33), Inches(1.6),
             title_html_runs, align=PP_ALIGN.CENTER, line_spacing=1.05)
    if sub:
        add_text(s, Inches(2), Inches(3.5), Inches(9.33), Inches(0.8),
                 sub, size=17, color=INK_3, align=PP_ALIGN.CENTER, line_spacing=1.4)
    return s


def slide_about():
    s = slide_kicker_section(
        "Who we are",
        [("One partner for the ", {"size": 48, "font": SERIF, "color": INK}),
         ("whole journey.", {"size": 48, "font": SERIF, "color": INK, "italic": True})],
    )

    # Two-column body
    # Left: copy
    add_text(s, Inches(0.8), Inches(4.6), Inches(7.0), Inches(2.2),
             "Incluzza is a one-stop, comprehensive and customizable solution for disability inclusion and accessibility. We are an offering of EnAble India, built around a practical framework that organizations can adopt as a whole or in individual “slices.”",
             size=16, color=INK_3, line_spacing=1.45)
    add_text(s, Inches(0.8), Inches(6.2), Inches(7.0), Inches(0.8),
             "You don’t have to solve everything at once. Start with the gap that matters most — and we’ll help you build from there.",
             size=18, italic=True, color=INK_2, font=SERIF, line_spacing=1.4)

    # Right: stats card
    card_x, card_y, card_w, card_h = Inches(8.3), Inches(4.4), Inches(4.2), Inches(2.5)
    add_rect(s, card_x, card_y, card_w, card_h, fill=PAPER, line=LINE, line_w=Pt(0.75),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)

    stats = [
        ("725+", "organizations supported"),
        ("291",  "roles opened for persons with disabilities"),
        ("22+",  "industries served"),
    ]
    for i, (num, lbl) in enumerate(stats):
        row_y = card_y + Inches(0.15) + Inches(0.78) * i
        if i > 0:
            add_rect(s, card_x + Inches(0.3), row_y - Inches(0.02),
                     card_w - Inches(0.6), Emu(9525), fill=LINE)
        add_text(s, card_x + Inches(0.3), row_y, Inches(1.6), Inches(0.6),
                 num, size=28, font=SERIF, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, card_x + Inches(1.95), row_y, card_w - Inches(2.2), Inches(0.6),
                 lbl, size=12, color=INK_3, anchor=MSO_ANCHOR.MIDDLE)


def slide_challenges():
    s = slide_kicker_section(
        "Where organizations get stuck",
        [("The challenges we hear ", {"size": 44, "font": SERIF, "color": INK}),
         ("most often.", {"size": 44, "font": SERIF, "color": INK, "italic": True})],
        "Sound familiar? You’re not alone — and these are exactly the barriers we help solve.",
    )

    items = [
        ("01", "“We want to hire, but don’t know how.”",        "Unclear pipelines, hesitant managers, and inaccessible interview processes."),
        ("02", "“Our people need confidence and awareness.”",   "Well-meaning teams that aren’t sure how to lead, support or include."),
        ("03", "“Our workplace and digital journeys have barriers.”", "Physical, digital and content gaps that exclude customers and colleagues."),
        ("04", "“Our policies and accountability need strengthening.”", "Good policies on paper that don’t translate into consistent practice."),
        ("05", "“We want products and services to work for more people.”", "Customer experiences that miss a large and loyal segment of the market."),
        ("06", "“Our language and content should be more inclusive.”",    "Messaging that inadvertently excludes the people it aims to reach."),
    ]

    cols, rows = 3, 2
    grid_x, grid_y = Inches(0.8), Inches(4.6)
    cell_w, cell_h = Inches(3.95), Inches(1.15)
    gap_x, gap_y = Inches(0.15), Inches(0.15)

    for i, (num, title, body) in enumerate(items):
        col = i % cols
        row = i // cols
        x = grid_x + (cell_w + gap_x) * col
        y = grid_y + (cell_h + gap_y) * row
        c = CAT[i % len(CAT)]

        # Card
        add_rect(s, x, y, cell_w, cell_h, fill=PAPER, line=LINE, line_w=Pt(0.75),
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # Top accent strip
        add_rect(s, x, y, cell_w, Inches(0.05), fill=c)
        # Number
        add_text(s, x + Inches(0.25), y + Inches(0.15), Inches(0.6), Inches(0.3),
                 num, size=12, font=SERIF, color=c)
        # Title
        add_text(s, x + Inches(0.25), y + Inches(0.4), cell_w - Inches(0.5), Inches(0.45),
                 title, size=13.5, font=SERIF, color=INK, line_spacing=1.15)
        # Body
        add_text(s, x + Inches(0.25), y + Inches(0.85), cell_w - Inches(0.5), Inches(0.3),
                 body, size=10.5, color=INK_3, line_spacing=1.25)


def slide_services():
    s = slide_kicker_section(
        "How we can support you",
        [("Eight connected areas. ", {"size": 44, "font": SERIF, "color": INK}),
         ("Pick what you need.", {"size": 44, "font": SERIF, "color": INK, "italic": True})],
        "Choose one slice to start, or combine several into a broader program. They reinforce each other.",
    )

    items = [
        ("01", "Accessibility",                "Physical, digital and content accessibility."),
        ("02", "Awareness & Sensitization",    "Workshops that build everyday confidence."),
        ("03", "Employment Support",           "Inclusive hiring, training and placement."),
        ("04", "Community Engagement",         "Interactive inclusion experiences."),
        ("05", "Accessible Products & Services","Make your customer journey work for everyone."),
        ("06", "Leadership & Commitment",      "Executive capability, ownership and accountability."),
        ("07", "Inclusive Policy",             "Policy design and review aligned with the RPwD Act."),
        ("08", "Inclusive Messaging",          "Language, content and communication that includes everyone."),
    ]

    cols, rows = 4, 2
    grid_x, grid_y = Inches(0.8), Inches(4.7)
    cell_w, cell_h = Inches(2.95), Inches(1.05)
    gap_x, gap_y = Inches(0.13), Inches(0.13)

    for i, (num, title, body) in enumerate(items):
        col = i % cols
        row = i // cols
        x = grid_x + (cell_w + gap_x) * col
        y = grid_y + (cell_h + gap_y) * row
        c = CAT[i % len(CAT)]

        add_rect(s, x, y, cell_w, cell_h, fill=PAPER, line=LINE, line_w=Pt(0.75),
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # Number pill
        pill = add_rect(s, x + Inches(0.22), y + Inches(0.15), Inches(0.45), Inches(0.28),
                        fill=SURFACE_2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = pill.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num
        r.font.name = SERIF; r.font.size = Pt(10); r.font.color.rgb = c; r.font.bold = True

        # Title
        add_text(s, x + Inches(0.22), y + Inches(0.42), cell_w - Inches(0.4), Inches(0.32),
                 title, size=13, font=SERIF, color=INK, bold=False, line_spacing=1.1)
        # Body
        add_text(s, x + Inches(0.22), y + Inches(0.7), cell_w - Inches(0.4), Inches(0.32),
                 body, size=9.5, color=INK_3, line_spacing=1.2)

    # Examples strip
    strip_y = Inches(7.05)
    add_text(s, Inches(0.8), strip_y, Inches(2.5), Inches(0.3),
             "EXAMPLES", size=10, bold=True, color=ACCENT, font=SANS)
    add_text(s, Inches(2.5), strip_y, Inches(10), Inches(0.3),
             "Awareness workshops · Sign-language training · Inclusive hiring programs · Accessibility audits · Sensitization modules · Policy advisory",
             size=11, color=INK_3)


def slide_value():
    s = slide_kicker_section(
        "What this means for your business",
        [("From an initiative to a ", {"size": 42, "font": SERIF, "color": INK}),
         ("measurable outcome.", {"size": 42, "font": SERIF, "color": INK, "italic": True})],
    )

    # From → To shift
    box_y = Inches(4.5)
    box_h = Inches(1.4)
    add_rect(s, Inches(0.8), box_y, Inches(11.73), box_h, fill=PAPER, line=LINE, line_w=Pt(0.75),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)

    # Left: Instead of asking
    add_text(s, Inches(1.0), box_y + Inches(0.2), Inches(4.6), Inches(0.3),
             "INSTEAD OF ASKING", size=10, bold=True, color=INK_3, font=SANS)
    add_text(s, Inches(1.0), box_y + Inches(0.5), Inches(4.6), Inches(0.9),
             "“Do we have a disability initiative?”",
             size=20, font=SERIF, italic=True, color=INK_2, line_spacing=1.2)

    # Arrow
    add_text(s, Inches(6.0), box_y, Inches(1.3), box_h,
             "→", size=28, color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Right: Imagine being able to say
    add_text(s, Inches(7.0), box_y + Inches(0.2), Inches(5.4), Inches(0.3),
             "IMAGINE BEING ABLE TO SAY", size=10, bold=True, color=ACCENT, font=SANS)
    add_text(s, Inches(7.0), box_y + Inches(0.5), Inches(5.4), Inches(0.9),
             "“People with disabilities can enter, contribute, grow and thrive here — and our products work for them too.”",
             size=18, font=SERIF, italic=True, color=INK, line_spacing=1.25)

    # 4 outcome tiles
    outcomes = [
        ("TALENT",   "Access a wider, more inclusive talent pool."),
        ("CULTURE",  "Build confidence, belonging and manager capability."),
        ("RISK",     "Strengthen accessibility and policy alignment."),
        ("BUSINESS", "Reach more customers with inclusive experiences."),
    ]
    tile_y = Inches(6.15)
    tile_w = Inches(2.85)
    tile_h = Inches(0.95)
    gap = Inches(0.12)
    for i, (tag, body) in enumerate(outcomes):
        x = Inches(0.8) + (tile_w + gap) * i
        add_rect(s, x, tile_y, tile_w, tile_h, fill=PAPER, line=LINE, line_w=Pt(0.75),
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # Subtle gradient feel via accent top
        add_rect(s, x, tile_y, tile_w, Inches(0.04), fill=ACCENT)
        add_text(s, x + Inches(0.2), tile_y + Inches(0.12), tile_w - Inches(0.4), Inches(0.25),
                 tag, size=10, bold=True, color=INK_3, font=SANS)
        add_text(s, x + Inches(0.2), tile_y + Inches(0.35), tile_w - Inches(0.4), Inches(0.55),
                 body, size=12, font=SERIF, color=INK, line_spacing=1.2)


def slide_why():
    s = slide_kicker_section(
        "Why Incluzza",
        [("A partner, ", {"size": 46, "font": SERIF, "color": INK}),
         ("not a vendor.", {"size": 46, "font": SERIF, "color": INK, "italic": True})],
    )

    items = [
        ("1", "Proven track record",       "725+ organizations supported, 291 job roles opened, across 22+ industries."),
        ("2", "End-to-end capability",     "Accessibility, talent, culture, leadership, policy, customer and messaging — under one roof."),
        ("3", "Customizable by design",    "Start with one priority. Scale into a full program as you grow."),
        ("4", "Practical delivery",        "Workshops, sign-language training, employment support and more — built for real workplaces."),
    ]

    cols = 2
    grid_x, grid_y = Inches(0.8), Inches(4.5)
    cell_w, cell_h = Inches(5.95), Inches(1.25)
    gap_x, gap_y = Inches(0.2), Inches(0.2)

    for i, (num, title, body) in enumerate(items):
        col = i % cols
        row = i // cols
        x = grid_x + (cell_w + gap_x) * col
        y = grid_y + (cell_h + gap_y) * row

        add_rect(s, x, y, cell_w, cell_h, fill=PAPER, line=LINE, line_w=Pt(0.75),
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # Big numeral
        add_text(s, x + Inches(0.3), y + Inches(0.15), Inches(1.2), Inches(1.0),
                 num, size=44, font=SERIF, color=ACCENT, line_spacing=1.0)
        # Title
        add_text(s, x + Inches(1.5), y + Inches(0.2), cell_w - Inches(1.7), Inches(0.4),
                 title, size=18, font=SERIF, color=INK)
        # Body
        add_text(s, x + Inches(1.5), y + Inches(0.65), cell_w - Inches(1.7), Inches(0.55),
                 body, size=12, color=INK_3, line_spacing=1.3)


def slide_roadmap_overview():
    s = slide_kicker_section(
        "Where to start",
        [("A practical path from ", {"size": 42, "font": SERIF, "color": INK}),
         ("where you are", {"size": 42, "font": SERIF, "color": INK, "italic": True}),
         (" to where you want to be.", {"size": 42, "font": SERIF, "color": INK})],
        "Every journey starts the same way: understand the starting point, then prioritize what matters most.",
    )

    phases = [
        ("01", "Discover",   "Understand your business, workforce and current employee experience."),
        ("02", "Diagnose",   "Identify barriers across people, process, technology and culture."),
        ("03", "Prioritize", "Rank what to act on by impact, urgency and feasibility."),
    ]

    # Rail line
    rail_y = Inches(5.05)
    add_rect(s, Inches(2.5), rail_y, Inches(8.3), Inches(0.025), fill=ACCENT)

    for i, (num, title, body) in enumerate(phases):
        x_center = Inches(2.0) + Inches(4.65) * i + Inches(2.33)
        # Circle
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  x_center - Inches(0.55), rail_y - Inches(0.55),
                                  Inches(1.1), Inches(1.1))
        circ.shadow.inherit = False
        circ.fill.solid(); circ.fill.fore_color.rgb = PAPER
        circ.line.color.rgb = INK; circ.line.width = Pt(1.25)
        add_text(s, x_center - Inches(0.55), rail_y - Inches(0.55), Inches(1.1), Inches(1.1),
                 num, size=18, font=SERIF, color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Title
        add_text(s, x_center - Inches(2.0), rail_y + Inches(0.85), Inches(4.0), Inches(0.45),
                 title, size=22, font=SERIF, color=INK, align=PP_ALIGN.CENTER)
        # Body
        add_text(s, x_center - Inches(2.0), rail_y + Inches(1.35), Inches(4.0), Inches(0.8),
                 body, size=13, color=INK_3, align=PP_ALIGN.CENTER, line_spacing=1.35)

    # Bottom strip
    add_rect(s, Inches(2.0), Inches(6.85), Inches(9.33), Inches(0.45),
             fill=PAPER, line=LINE, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_runs(s, Inches(2.2), Inches(6.85), Inches(9.0), Inches(0.45),
             [("From there:  ", {"size": 13, "color": INK_3}),
              ("Baseline → Priorities → Solution mix → Success measures.", {"size": 13, "color": INK, "bold": True})],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_workstreams():
    s = slide_kicker_section(
        "Key areas of focus",
        [("Five priorities, ", {"size": 44, "font": SERIF, "color": INK}),
         ("one owner each.", {"size": 44, "font": SERIF, "color": INK, "italic": True})],
        "A balanced program covers leadership, talent, access, culture and governance — together.",
    )

    items = [
        ("01", "Leadership",       ["Executive sponsor", "Clear ownership", "Accountability"]),
        ("02", "People & Talent",  ["Inclusive hiring", "Accessible onboarding", "Career development"]),
        ("03", "Accessibility",    ["Physical access", "Digital access", "Content & communication"]),
        ("04", "Culture",          ["Awareness", "Manager capability", "Employee networks"]),
        ("05", "Governance",       ["Policy", "Metrics", "Review cadence"]),
    ]

    cell_w, cell_h = Inches(2.32), Inches(2.3)
    gap = Inches(0.1)
    start_x = Inches(0.85)
    y = Inches(4.6)

    for i, (num, title, bullets) in enumerate(items):
        x = start_x + (cell_w + gap) * i
        add_rect(s, x, y, cell_w, cell_h, fill=PAPER, line=LINE, line_w=Pt(0.75),
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # Top accent strip
        add_rect(s, x, y, cell_w, Inches(0.05), fill=CAT[i % len(CAT)])
        # Number
        add_text(s, x + Inches(0.25), y + Inches(0.18), Inches(1), Inches(0.3),
                 num, size=12, font=SERIF, color=INK_3)
        # Title
        add_text(s, x + Inches(0.25), y + Inches(0.5), cell_w - Inches(0.5), Inches(0.4),
                 title, size=17, font=SERIF, color=INK, line_spacing=1.1)
        # Bullets
        for j, b in enumerate(bullets):
            add_text(s, x + Inches(0.25), y + Inches(1.0) + Inches(0.35) * j,
                     cell_w - Inches(0.5), Inches(0.3),
                     "• " + b, size=12, color=INK_3, line_spacing=1.2)


def slide_actions():
    s = slide_kicker_section(
        "A 12-month path",
        [("Quick wins early. ", {"size": 42, "font": SERIF, "color": INK}),
         ("Lasting change", {"size": 42, "font": SERIF, "color": INK, "italic": True}),
         (" over time.", {"size": 42, "font": SERIF, "color": INK})],
        "Each phase builds on the last — pairing visible progress with structural foundations.",
    )

    phases = [
        ("0 – 30 days",   "Understand", ["Baseline assessment", "Stakeholder interviews", "Accessibility & policy review"]),
        ("31 – 90 days",  "Build",      ["Leadership alignment", "Awareness workshops", "Inclusive hiring design"]),
        ("3 – 6 months",  "Scale",      ["Hiring & onboarding", "Accessibility improvements", "Employee engagement"]),
        ("6 – 12 months", "Sustain",    ["Metrics & governance", "Career pathways", "Continuous improvement"]),
    ]

    cell_w, cell_h = Inches(2.85), Inches(1.8)
    gap = Inches(0.15)
    start_x = Inches(0.85)
    y = Inches(4.55)

    for i, (when, phase, bullets) in enumerate(phases):
        x = start_x + (cell_w + gap) * i
        add_rect(s, x, y, cell_w, cell_h, fill=PAPER, line=LINE, line_w=Pt(0.75),
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # When
        add_text(s, x + Inches(0.25), y + Inches(0.18), cell_w - Inches(0.5), Inches(0.3),
                 when.upper(), size=10, bold=True, color=ACCENT, font=SANS)
        # Phase
        add_text(s, x + Inches(0.25), y + Inches(0.42), cell_w - Inches(0.5), Inches(0.45),
                 phase, size=22, font=SERIF, color=INK)
        # Bullets
        for j, b in enumerate(bullets):
            add_text(s, x + Inches(0.25), y + Inches(0.95) + Inches(0.26) * j,
                     cell_w - Inches(0.5), Inches(0.26),
                     "• " + b, size=11.5, color=INK_3, line_spacing=1.2)

    # Suggested actions row
    actions_y = Inches(6.55)
    actions = [
        ("LEADERSHIP", "Executive briefing & sponsor charter"),
        ("MANAGERS",   "Inclusive hiring & workplace confidence"),
        ("TALENT",     "Accessible recruitment & onboarding"),
        ("ACCESS",     "Physical, digital & content fixes"),
        ("CULTURE",    "Awareness & lived-experience learning"),
        ("GOVERNANCE", "Policy refresh, dashboard, quarterly review"),
    ]
    add_text(s, Inches(0.8), actions_y, Inches(2), Inches(0.3),
             "SUGGESTED ACTIONS", size=10, bold=True, color=ACCENT, font=SANS)
    add_text(s, Inches(2.7), actions_y, Inches(10.5), Inches(0.5),
             "  ·  ".join(a[1] for a in actions),
             size=11, color=INK_3, line_spacing=1.3)


def slide_measure():
    s = slide_kicker_section(
        "How we measure progress",
        [("Beyond activity metrics — to ", {"size": 40, "font": SERIF, "color": INK}),
         ("real outcomes.", {"size": 40, "font": SERIF, "color": INK, "italic": True})],
        "A four-level model that tracks investment, delivery, behavior change and business impact.",
    )

    levels = [
        ("1", "Input",   "What we invest",        ["Training hours", "Resources", "Owners assigned"]),
        ("2", "Output",  "What we deliver",       ["People reached", "Candidates supported", "Issues addressed"]),
        ("3", "Outcome", "What changes",          ["Representation", "Retention", "Inclusion confidence"]),
        ("4", "Impact",  "What the business gains",["Talent access", "Customer experience", "Accessibility maturity"]),
    ]

    cell_w, cell_h = Inches(2.95), Inches(2.1)
    gap = Inches(0.13)
    start_x = Inches(0.8)
    y = Inches(4.55)

    for i, (num, title, q, bullets) in enumerate(levels):
        x = start_x + (cell_w + gap) * i
        add_rect(s, x, y, cell_w, cell_h, fill=PAPER, line=LINE, line_w=Pt(0.75),
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # Big number
        add_text(s, x + Inches(0.3), y + Inches(0.15), Inches(0.8), Inches(0.55),
                 num, size=30, font=SERIF, color=ACCENT)
        # Title
        add_text(s, x + Inches(0.3), y + Inches(0.7), cell_w - Inches(0.6), Inches(0.4),
                 title, size=22, font=SERIF, color=INK)
        # Question
        add_text(s, x + Inches(0.3), y + Inches(1.1), cell_w - Inches(0.6), Inches(0.3),
                 q, size=12, italic=True, color=INK_3, font=SERIF)
        # Bullets
        for j, b in enumerate(bullets):
            add_text(s, x + Inches(0.3), y + Inches(1.4) + Inches(0.22) * j,
                     cell_w - Inches(0.6), Inches(0.22),
                     "• " + b, size=11, color=INK_3, line_spacing=1.2)

    # Principle strip
    add_rect(s, Inches(0.8), Inches(6.85), Inches(11.73), Inches(0.45),
             fill=PAPER, line=LINE, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_runs(s, Inches(1.0), Inches(6.85), Inches(11.4), Inches(0.45),
             [("OUR PRINCIPLE   ", {"size": 11, "bold": True, "color": ACCENT, "font": SANS}),
              ("We establish the baseline first. Targets are set together, after the data — never invented in advance.",
               {"size": 13, "color": INK_2, "font": SANS})],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_close():
    s = prs.slides.add_slide(BLANK)
    slide_background(s, SURFACE_2)
    add_ornament_orbs(s)

    # Kicker
    kicker_pill(s, Inches(5.37), Inches(0.9), "Next step")

    # Title
    add_runs(s, Inches(1), Inches(1.4), Inches(11.33), Inches(1.2),
             [("Let’s start with a ", {"size": 52, "font": SERIF, "color": INK}),
              ("conversation.", {"size": 52, "font": SERIF, "color": INK, "italic": True})],
             align=PP_ALIGN.CENTER, line_spacing=1.05)
    add_text(s, Inches(2), Inches(2.55), Inches(9.33), Inches(0.5),
             "A short discovery session is the best way to understand your starting point and identify the highest-impact place to begin.",
             size=15, color=INK_3, align=PP_ALIGN.CENTER)

    # 3 step cards
    steps = [
        ("FIRST",   "A discovery session",   "We listen, learn about your organization, and identify the top three barriers to focus on."),
        ("THEN",    "A tailored proposal",   "A short, prioritized roadmap with owners, actions and a clear path to measurable outcomes."),
        ("FINALLY", "Delivery & review",     "We work alongside your team, with quarterly reviews to keep the program on track."),
    ]

    cell_w, cell_h = Inches(3.95), Inches(2.0)
    gap = Inches(0.2)
    start_x = Inches(0.8)
    y = Inches(3.5)

    for i, (when, title, body) in enumerate(steps):
        x = start_x + (cell_w + gap) * i
        add_rect(s, x, y, cell_w, cell_h, fill=PAPER, line=LINE, line_w=Pt(0.75),
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # Top accent
        add_rect(s, x, y, cell_w, Inches(0.05), fill=ACCENT)
        add_text(s, x + Inches(0.3), y + Inches(0.2), cell_w - Inches(0.6), Inches(0.3),
                 when, size=11, bold=True, color=ACCENT, font=SANS)
        add_text(s, x + Inches(0.3), y + Inches(0.5), cell_w - Inches(0.6), Inches(0.5),
                 title, size=22, font=SERIF, color=INK, line_spacing=1.1)
        add_text(s, x + Inches(0.3), y + Inches(1.05), cell_w - Inches(0.6), Inches(0.85),
                 body, size=12, color=INK_3, line_spacing=1.35)

    # CTA + footer
    add_text(s, Inches(1), Inches(6.0), Inches(11.33), Inches(0.5),
             "Get in touch →",
             size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(6.6), Inches(11.33), Inches(0.4),
             "Thank you for your time. We’d love to continue the conversation.",
             size=14, italic=True, color=INK_3, font=SERIF, align=PP_ALIGN.CENTER)

    # Brand strip at bottom
    add_rect(s, 0, Inches(7.15), SW, Inches(0.35), fill=SURFACE_2)
    add_text(s, Inches(0.5), Inches(7.18), Inches(12.33), Inches(0.3),
             "Incluzza  ·  An offering of EnAble India  ·  A partner for disability inclusion",
             size=11, color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── Build the deck ────────────────────────────────────────────────
slide_cover()
slide_part_marker("Part 1", "The client pitch", "Who we are, the challenges we solve, and the value we bring.")
slide_about()
slide_challenges()
slide_services()
slide_value()
slide_why()
slide_part_marker("Part 2", "Your roadmap", "Where to start, what to focus on, what to do, and how to measure progress.")
slide_roadmap_overview()
slide_workstreams()
slide_actions()
slide_measure()
slide_close()

OUT = "/Users/harshithasrinivas/Desktop/incluzza-site/Incluzza_Client_Pitch_and_Roadmap.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Total slides: {len(prs.slides)}")